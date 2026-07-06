#!/usr/bin/env python3
"""
Reads pack_content.json and produces the six Oolio GTM pack files alongside
it. Equivalent to scripts/build_pack.js but written for environments where
the Node toolchain is unavailable.

Usage:
    python3 build_pack.py <path to pack_content.json>

Outputs (in the JSON's directory):
    01_One_Pager_<v>.pptx
    02_Supporting_Deck_<v>.pptx
    03_Sales_Playbook_<v>.docx
    04_Account_Management_Playbook_<v>.docx
    05_Onboarding_Playbook_<v>.docx
    06_Marketing_Pack_<v>.docx
"""

import json
import os
import sys
from datetime import datetime

from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsmap
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu as PEmu
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


# ---- Brand ----
PURPLE = RGBColor(0x67, 0x3A, 0xB6)
PURPLE_DARK = RGBColor(0x5E, 0x35, 0xB1)
PURPLE_LIGHT_HEX = "F9F8FC"
HEADER_FILL_HEX = "EDE7F6"
ZEBRA_HEX = "F5F1FA"
RULE_HEX = "D5D0E1"
TEXT_DARK = RGBColor(0x1F, 0x1A, 0x2E)
TEXT_GREY = RGBColor(0x5A, 0x55, 0x66)
TEXT_MUTE = RGBColor(0x8C, 0x87, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PURPLE_HEX = "673AB6"
PURPLE_DARK_HEX = "5E35B1"
TEXT_DARK_HEX = "1F1A2E"
TEXT_GREY_HEX = "5A5566"

PPTX_PURPLE = PRGB(0x67, 0x3A, 0xB6)
PPTX_PURPLE_DARK = PRGB(0x5E, 0x35, 0xB1)
PPTX_PURPLE_LIGHT = PRGB(0xF9, 0xF8, 0xFC)
PPTX_TEXT_DARK = PRGB(0x1F, 0x1A, 0x2E)
PPTX_TEXT_GREY = PRGB(0x5A, 0x55, 0x66)
PPTX_TEXT_MUTE = PRGB(0x8C, 0x87, 0x9A)
PPTX_WHITE = PRGB(0xFF, 0xFF, 0xFF)
PPTX_ACCENT = PRGB(0xED, 0xE7, 0xF6)


# ============================================================
# Helpers
# ============================================================

def show(v, fallback="[TBC]"):
    if v is None:
        return fallback
    s = str(v)
    if s.strip() == "" or s == "[GAP]":
        return "[NOT YET FILLED]"
    return s


def set_cell_shading(cell, hex_color):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tc_pr.append(shd)


def set_cell_borders(cell, color_hex=RULE_HEX, sz="4"):
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_borders = OxmlElement("w:tcBorders")
    for edge in ("top", "left", "bottom", "right"):
        b = OxmlElement(f"w:{edge}")
        b.set(qn("w:val"), "single")
        b.set(qn("w:sz"), sz)
        b.set(qn("w:color"), color_hex)
        tc_borders.append(b)
    tc_pr.append(tc_borders)


def add_run(paragraph, text, *, bold=False, italic=False, color=None, size=11, font="Arial"):
    run = paragraph.add_run(text)
    run.font.name = font
    run.bold = bold
    run.italic = italic
    if color is not None:
        run.font.color.rgb = color
    run.font.size = Pt(size)
    return run


def para(doc_or_cell, text, *, bold=False, italic=False, color=None, size=11, alignment=None, after=4, before=0, font="Arial"):
    p = doc_or_cell.add_paragraph()
    if alignment is not None:
        p.alignment = alignment
    p.paragraph_format.space_after = Pt(after)
    p.paragraph_format.space_before = Pt(before)
    add_run(p, text, bold=bold, italic=italic, color=color, size=size, font=font)
    return p


def heading(doc, text, *, level=1):
    sizes = {1: 16, 2: 13, 3: 11}
    size = sizes.get(level, 13)
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18 if level == 1 else 14)
    p.paragraph_format.space_after = Pt(9 if level == 1 else 6)
    add_run(p, text, bold=True, color=PURPLE_DARK, size=size)
    if level == 1:
        # Underline rule via paragraph border
        p_pr = p._p.get_or_add_pPr()
        pbdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:color"), PURPLE_HEX)
        bottom.set(qn("w:space"), "1")
        pbdr.append(bottom)
        p_pr.append(pbdr)
    return p


def bullet(doc, text):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.space_after = Pt(3)
    for r in p.runs:
        r.font.name = "Arial"
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT_DARK
    if not p.runs:
        add_run(p, text, size=11)
    else:
        p.runs[0].text = text
    return p


def numbered(doc, text):
    p = doc.add_paragraph(style="List Number")
    p.paragraph_format.space_after = Pt(3)
    for r in p.runs:
        r.font.name = "Arial"
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT_DARK
    if not p.runs:
        add_run(p, text, size=11)
    else:
        p.runs[0].text = text
    return p


def header_row(table, labels, *, widths_in=None):
    row = table.rows[0]
    for i, label in enumerate(labels):
        cell = row.cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        p.paragraph_format.space_after = Pt(2)
        add_run(p, label, bold=True, color=PURPLE_DARK, size=10)
        set_cell_shading(cell, HEADER_FILL_HEX)
        set_cell_borders(cell)
        cell.vertical_alignment = WD_ALIGN_VERTICAL.TOP
    if widths_in:
        for i, w in enumerate(widths_in):
            for r in table.rows:
                r.cells[i].width = Inches(w)


def add_data_row(table, cells, *, zebra=False, bold_first=False):
    row = table.add_row()
    for i, val in enumerate(cells):
        cell = row.cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        p.paragraph_format.space_after = Pt(2)
        add_run(p, str(val), bold=(bold_first and i == 0), size=10)
        if zebra:
            set_cell_shading(cell, ZEBRA_HEX)
        set_cell_borders(cell)
        cell.vertical_alignment = WD_ALIGN_VERTICAL.TOP


def header_footer(doc, kind, product):
    section = doc.sections[0]
    header = section.header
    h_para = header.paragraphs[0]
    h_para.text = ""
    add_run(h_para, f"Oolio  /  {kind}  /  {show(product.get('name'))}",
            bold=True, color=PURPLE, size=8)
    add_run(h_para, f"\t{show(product.get('version'), 'v0.1')}  /  {show(product.get('last_reviewed'))}",
            color=TEXT_GREY, size=8)
    h_para.paragraph_format.tab_stops.add_tab_stop(Inches(6.5), WD_ALIGN_PARAGRAPH.RIGHT)
    p_pr = h_para._p.get_or_add_pPr()
    pbdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")
    bottom.set(qn("w:color"), PURPLE_HEX)
    bottom.set(qn("w:space"), "4")
    pbdr.append(bottom)
    p_pr.append(pbdr)

    footer = section.footer
    f_para = footer.paragraphs[0]
    f_para.text = ""
    add_run(f_para, "Oolio GTM & Product Enablement", color=TEXT_GREY, size=8)
    add_run(f_para, "\t", color=TEXT_GREY, size=8)
    add_run(f_para, "Page ", color=TEXT_GREY, size=8)
    # Page number field
    page_run = f_para.add_run()
    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.text = "PAGE"
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "end")
    page_run._r.append(fld_char1)
    page_run._r.append(instr)
    page_run._r.append(fld_char2)
    page_run.font.color.rgb = TEXT_GREY
    page_run.font.size = Pt(8)
    f_para.paragraph_format.tab_stops.add_tab_stop(Inches(6.5), WD_ALIGN_PARAGRAPH.RIGHT)


def metadata_table(doc, product):
    table = doc.add_table(rows=0, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    rows = [
        ("Product", show(product.get("name"))),
        ("Owner (PM)", show((product.get("owner_pm") or {}).get("name"))),
        ("GTM partner", show((product.get("gtm_partner") or {}).get("name"))),
        ("Version", show(product.get("version"), "v0.1")),
        ("Last reviewed", show(product.get("last_reviewed"))),
        ("Anchor Jira epic", show(product.get("anchor_jira_epic"))),
    ]
    for label, val in rows:
        r = table.add_row()
        c0, c1 = r.cells
        c0.text = ""
        c1.text = ""
        p0 = c0.paragraphs[0]
        add_run(p0, label, bold=True, color=PURPLE_DARK, size=10)
        set_cell_shading(c0, HEADER_FILL_HEX)
        set_cell_borders(c0)
        p1 = c1.paragraphs[0]
        add_run(p1, val, size=10)
        set_cell_borders(c1)
        c0.width = Inches(1.6)
        c1.width = Inches(4.9)


def title_block(doc, title, subtitle, product):
    p1 = doc.add_paragraph()
    p1.paragraph_format.space_after = Pt(2)
    add_run(p1, "OOLIO  /  GTM & PRODUCT ENABLEMENT", bold=True, color=PURPLE, size=8)

    p2 = doc.add_paragraph()
    p2.paragraph_format.space_after = Pt(2)
    add_run(p2, title, bold=True, size=24)

    p3 = doc.add_paragraph()
    p3.paragraph_format.space_after = Pt(2)
    add_run(p3, subtitle, italic=True, color=TEXT_GREY, size=12)

    p4 = doc.add_paragraph()
    p4.paragraph_format.space_after = Pt(14)
    add_run(p4, f"For {show(product.get('name'))}  /  {show(product.get('version'), 'v0.1')}",
            bold=True, color=PURPLE_DARK, size=11)

    rule = doc.add_paragraph()
    rule.paragraph_format.space_before = Pt(2)
    rule.paragraph_format.space_after = Pt(8)
    p_pr = rule._p.get_or_add_pPr()
    pbdr = OxmlElement("w:pBdr")
    top = OxmlElement("w:top")
    top.set(qn("w:val"), "single")
    top.set(qn("w:sz"), "6")
    top.set(qn("w:color"), PURPLE_HEX)
    top.set(qn("w:space"), "1")
    pbdr.append(top)
    p_pr.append(pbdr)

    heading(doc, "Document metadata", level=2)
    metadata_table(doc, product)
    page_break(doc)


def page_break(doc):
    p = doc.add_paragraph()
    r = p.add_run()
    r.add_break(WD_BREAK.PAGE)


# ============================================================
# 03. Sales Playbook
# ============================================================

def build_sales(content, product, out_path):
    s = content.get("sales") or {}
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.9)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)
    style.font.color.rgb = TEXT_DARK

    header_footer(doc, "Sales Playbook", product)
    title_block(doc, "Sales Playbook",
                "How to qualify, demo, handle objections, and close.", product)

    heading(doc, "1. Product summary", level=1)
    para(doc, show(s.get("summary")))

    heading(doc, "2. Discovery", level=1)
    heading(doc, "2.1  Discovery questions", level=2)
    for q in (s.get("discovery_questions") or []):
        numbered(doc, show(q))

    heading(doc, "2.2  Qualification criteria", level=2)
    qual = s.get("qualification") or []
    table = doc.add_table(rows=1, cols=3)
    header_row(table, ["Criterion", "Qualified in", "Qualified out"], widths_in=[1.6, 2.45, 2.45])
    for i, q in enumerate(qual):
        add_data_row(table, [show(q.get("criterion")), show(q.get("qualified_in")), show(q.get("qualified_out"))],
                     zebra=(i % 2 == 1), bold_first=True)

    heading(doc, "3. Demo flow and talk track", level=1)
    for i, step in enumerate(s.get("demo_steps") or [], start=1):
        heading(doc, f"Step {i}. {show(step.get('step'))} ({step.get('minutes', '[TBC]')} min)", level=3)
        p = doc.add_paragraph()
        add_run(p, "Show: ", bold=True, italic=True, color=TEXT_GREY, size=11)
        add_run(p, show(step.get("what_to_show")), italic=True, color=TEXT_GREY, size=11)
        p2 = doc.add_paragraph()
        add_run(p2, "Say: ", bold=True, size=11)
        add_run(p2, show(step.get("talk_track")), size=11)

    heading(doc, "4. Objection handling", level=1)
    obj_table = doc.add_table(rows=1, cols=2)
    header_row(obj_table, ["Objection", "Response and proof"], widths_in=[2.1, 4.4])
    for i, o in enumerate(s.get("objections") or []):
        add_data_row(obj_table, [show(o.get("objection")), show(o.get("response"))],
                     zebra=(i % 2 == 1), bold_first=True)

    heading(doc, "5. Pricing and packaging quick reference", level=1)
    heading(doc, "5.1  Packages", level=3)
    for p_pkg in ((s.get("pricing") or {}).get("packages") or []):
        para(doc, f"{show(p_pkg.get('name'))}. For {show(p_pkg.get('for'))}. Includes {show(p_pkg.get('includes'))}. {show(p_pkg.get('headline_price'))}.")
    heading(doc, "5.2  Discounting authority", level=3)
    para(doc, show((s.get("pricing") or {}).get("discount_authority")))
    heading(doc, "5.3  Source of truth", level=3)
    para(doc, show((s.get("pricing") or {}).get("source_of_truth")))

    heading(doc, "6. Close and handover", level=1)
    heading(doc, "6.1  Close motion", level=3)
    para(doc, show(s.get("close_motion")))
    heading(doc, "6.2  Handover checklist", level=3)
    for item in (s.get("handover_checklist") or []):
        numbered(doc, show(item))

    heading(doc, "7. Appendix", level=1)
    bullet(doc, f"Linked Supporting Deck: 02_Supporting_Deck_{show(product.get('version'))}.pptx")
    bullet(doc, f"Linked One-Pager: 01_One_Pager_{show(product.get('version'))}.pptx")
    bullet(doc, f"Linked Marketing Pack: 06_Marketing_Pack_{show(product.get('version'))}.docx")

    doc.save(out_path)


# ============================================================
# 04. Account Management Playbook
# ============================================================

def build_am(content, product, out_path):
    a = content.get("am") or {}
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.9)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)
    style.font.color.rgb = TEXT_DARK

    header_footer(doc, "Account Management Playbook", product)
    title_block(doc, "Account Management Playbook",
                "Adoption, expansion, renewal, and at-risk plays.", product)

    heading(doc, "1. Product summary for AMs", level=1)
    para(doc, show(a.get("summary")))
    p = doc.add_paragraph()
    add_run(p, f"At 90 days post go-live: {show(a.get('healthy_at_90d'))}",
            italic=True, color=TEXT_GREY, size=11)

    heading(doc, "2. Adoption signals and red flags", level=1)
    table = doc.add_table(rows=1, cols=3)
    header_row(table, ["Feature", "Healthy", "Red flag"], widths_in=[1.8, 2.35, 2.35])
    for i, row in enumerate(a.get("adoption_signals") or []):
        add_data_row(table, [show(row.get("feature")), show(row.get("healthy")), show(row.get("red_flag"))],
                     zebra=(i % 2 == 1), bold_first=True)

    heading(doc, "3. Expansion triggers and upsell pathways", level=1)
    et = doc.add_table(rows=1, cols=3)
    header_row(et, ["Trigger", "Suggested next product", "Talk track"], widths_in=[2.1, 2.1, 2.3])
    for i, t in enumerate(a.get("expansion_triggers") or []):
        add_data_row(et, [show(t.get("trigger")), show(t.get("next_product")), show(t.get("talk_track"))],
                     zebra=(i % 2 == 1))

    heading(doc, "4. Renewal play", level=1)
    heading(doc, "4.1  Timeline", level=3)
    for t in (a.get("renewal_timeline") or []):
        numbered(doc, f"{show(t.get('checkpoint'))}: {show(t.get('action'))}")
    heading(doc, "4.2  Required artifacts", level=3)
    for it in (a.get("renewal_artifacts") or []):
        bullet(doc, show(it))

    heading(doc, "5. At-risk play", level=1)
    heading(doc, "5.1  Trigger criteria", level=3)
    for it in (a.get("at_risk_triggers") or []):
        bullet(doc, show(it))
    heading(doc, "5.2  Recovery sequence", level=3)
    for st in (a.get("recovery_sequence") or []):
        numbered(doc, f"{show(st.get('step'))}: {show(st.get('action'))}")

    heading(doc, "6. Quarterly Business Review template", level=1)
    for s_item in (a.get("qbr_sections") or []):
        numbered(doc, show(s_item))

    heading(doc, "7. Appendix", level=1)
    bullet(doc, f"Linked Sales Playbook: 03_Sales_Playbook_{show(product.get('version'))}.docx")
    bullet(doc, f"Linked Onboarding Playbook: 05_Onboarding_Playbook_{show(product.get('version'))}.docx")

    doc.save(out_path)


# ============================================================
# 05. Onboarding Playbook
# ============================================================

def build_onboarding(content, product, out_path):
    o = content.get("onboarding") or {}
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.9)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)
    style.font.color.rgb = TEXT_DARK

    header_footer(doc, "Onboarding Playbook", product)
    title_block(doc, "Onboarding Playbook",
                "From contract signed to first-week health check.", product)

    heading(doc, "1. Definition of done", level=1)
    for item in (o.get("definition_of_done") or []):
        bullet(doc, show(item))

    heading(doc, "2. Configuration checklist", level=1)
    cc = doc.add_table(rows=1, cols=3)
    header_row(cc, ["Step", "Owner", "Due"], widths_in=[3.5, 1.5, 1.5])
    for i, s_item in enumerate(o.get("config_checklist") or []):
        add_data_row(cc, [show(s_item.get("step")), show(s_item.get("owner")), show(s_item.get("due"))],
                     zebra=(i % 2 == 1))

    heading(doc, "3. Data migration plan", level=1)
    heading(doc, "3.1  In-scope data", level=3)
    mig = (o.get("migration") or {}).get("in_scope") or []
    mt = doc.add_table(rows=1, cols=4)
    header_row(mt, ["Domain", "Source", "Volume", "Owner"], widths_in=[1.7, 2.1, 1.3, 1.4])
    for i, m in enumerate(mig):
        add_data_row(mt, [show(m.get("domain")), show(m.get("source")), show(m.get("record_count")), show(m.get("owner"))],
                     zebra=(i % 2 == 1), bold_first=True)
    heading(doc, "3.2  Mapping decisions", level=3)
    para(doc, show((o.get("migration") or {}).get("mapping_decisions")))
    heading(doc, "3.3  Validation steps", level=3)
    for it in ((o.get("migration") or {}).get("validation") or []):
        numbered(doc, show(it))

    heading(doc, "4. Staff training plan", level=1)
    heading(doc, "4.1  Manager training", level=3)
    for it in ((o.get("training") or {}).get("manager") or []):
        bullet(doc, show(it))
    heading(doc, "4.2  Frontline training", level=3)
    for it in ((o.get("training") or {}).get("frontline") or []):
        bullet(doc, show(it))
    heading(doc, "4.3  Materials provided", level=3)
    for it in ((o.get("training") or {}).get("materials") or []):
        bullet(doc, show(it))

    heading(doc, "5. Go-live validation", level=1)
    heading(doc, "5.1  T-24h pre-flight", level=3)
    for it in ((o.get("golive") or {}).get("preflight") or []):
        bullet(doc, show(it))
    heading(doc, "5.2  T-0 cutover", level=3)
    co = doc.add_table(rows=1, cols=3)
    header_row(co, ["Step", "Owner", "ETA"], widths_in=[3.5, 1.7, 1.3])
    for i, st in enumerate((o.get("golive") or {}).get("cutover") or []):
        add_data_row(co, [show(st.get("step")), show(st.get("owner")), show(st.get("eta"))],
                     zebra=(i % 2 == 1))
    heading(doc, "5.3  T+24h post-live", level=3)
    for it in ((o.get("golive") or {}).get("postlive") or []):
        bullet(doc, show(it))

    heading(doc, "6. First-week health check", level=1)
    hc = doc.add_table(rows=1, cols=3)
    header_row(hc, ["Check", "Pass criteria", "Fail action"], widths_in=[1.7, 2.4, 2.4])
    for i, h in enumerate(o.get("healthcheck") or []):
        add_data_row(hc, [show(h.get("check")), show(h.get("pass")), show(h.get("fail_action"))],
                     zebra=(i % 2 == 1), bold_first=True)

    heading(doc, "7. Handover to Account Management", level=1)
    for it in (o.get("handover") or []):
        numbered(doc, show(it))

    doc.save(out_path)


# ============================================================
# 06. Marketing Pack
# ============================================================

def build_marketing(content, product, out_path):
    m = content.get("marketing") or {}
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.9)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)
    style.font.color.rgb = TEXT_DARK

    header_footer(doc, "Marketing Pack", product)
    title_block(doc, "Marketing Pack",
                "Launch, social, email, enablement, and campaign brief.", product)

    tier = m.get("tier", 2)
    heading(doc, "1. Pack overview", level=1)
    para(doc, f"Tier {tier} pack.")
    if tier == 1:
        para(doc, "Tier 1 includes the full set of artifacts: launch announcement, social, prospect and customer email sequences, sales enablement note, and a paid campaign brief.", italic=True, color=TEXT_GREY)
    elif tier == 2:
        para(doc, "Tier 2 includes the launch announcement, social, prospect and customer email sequences, and the sales enablement note. No paid campaign brief.", italic=True, color=TEXT_GREY)
    else:
        para(doc, "Tier 3 is customer-facing only: customer email sequence and the sales enablement note.", italic=True, color=TEXT_GREY)

    if tier in (1, 2):
        heading(doc, "2. Launch announcement", level=1)
        heading(doc, "2.1  LinkedIn long-form", level=3)
        for ln in (m.get("launch") or {}).get("linkedin_long", "[TBC]").split("\n"):
            if ln.strip():
                para(doc, ln)
            else:
                doc.add_paragraph()
        heading(doc, "2.2  LinkedIn short-form", level=3)
        para(doc, show((m.get("launch") or {}).get("linkedin_short")))
        heading(doc, "2.3  Partner channel", level=3)
        para(doc, show((m.get("launch") or {}).get("partner_channel")))

        heading(doc, "3. Social posts", level=1)
        for i, p in enumerate(m.get("social_posts") or [], start=1):
            heading(doc, f"Post {i}", level=3)
            para(doc, show(p.get("body")))
            p_v = doc.add_paragraph()
            add_run(p_v, "Visual: ", bold=True, italic=True, color=TEXT_GREY, size=10)
            add_run(p_v, show(p.get("visual")), italic=True, color=TEXT_GREY, size=10)

        heading(doc, "4. Email sequences", level=1)
        heading(doc, "4.1  Prospect sequence", level=2)
        for i, e in enumerate((m.get("emails") or {}).get("prospect") or [], start=1):
            heading(doc, f"Email {i}", level=3)
            ph = doc.add_paragraph()
            add_run(ph, "Subject: ", bold=True, size=11)
            add_run(ph, show(e.get("subject")), size=11)
            pph = doc.add_paragraph()
            add_run(pph, "Pre-header: ", bold=True, size=11)
            add_run(pph, show(e.get("preheader")), size=11)
            para(doc, show(e.get("body")))
            pcta = doc.add_paragraph()
            add_run(pcta, "CTA: ", bold=True, color=PURPLE_DARK, size=11)
            add_run(pcta, show(e.get("cta")), color=PURPLE_DARK, size=11)

    heading(doc, "5. Customer email sequence", level=1)
    for i, e in enumerate((m.get("emails") or {}).get("customer") or [], start=1):
        heading(doc, f"Email {i}", level=3)
        ph = doc.add_paragraph()
        add_run(ph, "Subject: ", bold=True, size=11)
        add_run(ph, show(e.get("subject")), size=11)
        pph = doc.add_paragraph()
        add_run(pph, "Pre-header: ", bold=True, size=11)
        add_run(pph, show(e.get("preheader")), size=11)
        para(doc, show(e.get("body")))
        pcta = doc.add_paragraph()
        add_run(pcta, "CTA: ", bold=True, color=PURPLE_DARK, size=11)
        add_run(pcta, show(e.get("cta")), color=PURPLE_DARK, size=11)

    heading(doc, "6. Sales enablement note", level=1)
    para(doc, show(m.get("enablement_note")))

    if tier == 1:
        heading(doc, "7. Campaign brief", level=1)
        cb = m.get("campaign_brief") or {}
        rows = [
            ("Name", show(cb.get("name"))),
            ("Objective", show(cb.get("objective"))),
            ("Audience", show(cb.get("audience"))),
            ("Headline message", show(cb.get("headline_message"))),
            ("Channels", ", ".join(cb.get("channels") or ["[TBC]"])),
            ("Budget", show(cb.get("budget"))),
            ("Run dates", show(cb.get("run_dates"))),
            ("Success metric", show(cb.get("success_metric"))),
            ("Owner", show(cb.get("owner"))),
        ]
        t = doc.add_table(rows=0, cols=2)
        for label, val in rows:
            r = t.add_row()
            c0, c1 = r.cells
            c0.text = ""
            c1.text = ""
            p0 = c0.paragraphs[0]
            add_run(p0, label, bold=True, color=PURPLE_DARK, size=10)
            set_cell_shading(c0, HEADER_FILL_HEX)
            set_cell_borders(c0)
            p1 = c1.paragraphs[0]
            add_run(p1, val, size=10)
            set_cell_borders(c1)
            c0.width = Inches(1.6)
            c1.width = Inches(4.9)

    doc.save(out_path)


# ============================================================
# PPTX helpers
# ============================================================

def add_text_box(slide, x, y, w, h, text, *, size=18, bold=False, color=None,
                 font="Arial", anchor=None, align=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.margin_left = PEmu(0)
    tf.margin_right = PEmu(0)
    tf.margin_top = PEmu(0)
    tf.margin_bottom = PEmu(0)

    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = PPt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb


def add_filled_rect(slide, x, y, w, h, fill_color, *, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_card(slide, x, y, w, h, fill_color, line_color=None):
    return add_filled_rect(slide, x, y, w, h, fill_color, line_color=line_color)


# ============================================================
# 01. One-Pager (single slide, A4 portrait)
# ============================================================

def build_one_pager(content, product, out_path):
    op = content.get("one_pager") or {}
    prs = Presentation()
    # A4 portrait: 8.27in x 11.69in
    prs.slide_width = PInches(8.27)
    prs.slide_height = PInches(11.69)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, PPTX_PURPLE_LIGHT)

    # Top purple bar
    add_filled_rect(slide, 0, 0, prs.slide_width, PInches(0.9), PPTX_PURPLE_DARK)
    add_text_box(slide, PInches(0.4), PInches(0.18), PInches(7.5), PInches(0.3),
                 "OOLIO  /  GTM ONE-PAGER", size=10, bold=True, color=PPTX_WHITE)
    add_text_box(slide, PInches(0.4), PInches(0.4), PInches(7.5), PInches(0.55),
                 show(product.get("name")), size=22, bold=True, color=PPTX_WHITE)

    # Tagline band
    add_filled_rect(slide, PInches(0.4), PInches(1.05), PInches(7.47), PInches(0.7), PPTX_ACCENT)
    add_text_box(slide, PInches(0.6), PInches(1.18), PInches(7.07), PInches(0.5),
                 show(op.get("tagline")), size=16, bold=True, color=PPTX_PURPLE_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Four blocks
    block_y = PInches(1.95)
    block_h = PInches(2.3)
    block_w = PInches(3.7)
    gutter = PInches(0.07)

    blocks = [
        ("WHAT IS IT", show(op.get("what")), 0, 0),
        ("WHO IS IT FOR", show(op.get("who")), 0, 1),
        ("THE PROBLEM", show(op.get("problem")), 1, 0),
        ("THE OUTCOME", show(op.get("outcome")), 1, 1),
    ]
    for label, val, row, col in blocks:
        x = PInches(0.4 + col * 3.77)
        y = PInches(1.95 + row * 2.4)
        add_filled_rect(slide, x, y, block_w, block_h, PPTX_WHITE, line_color=PPTX_PURPLE)
        add_text_box(slide, x + PInches(0.2), y + PInches(0.15), block_w - PInches(0.4), PInches(0.3),
                     label, size=9, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(slide, x + PInches(0.2), y + PInches(0.5), block_w - PInches(0.4), block_h - PInches(0.7),
                     val, size=11, color=PPTX_TEXT_DARK)

    # Proof point band
    proof_y = PInches(6.85)
    add_filled_rect(slide, PInches(0.4), proof_y, PInches(7.47), PInches(1.4), PPTX_PURPLE_DARK)
    add_text_box(slide, PInches(0.6), proof_y + PInches(0.15), PInches(7.07), PInches(0.3),
                 "PROOF POINT", size=10, bold=True, color=PPTX_WHITE)
    add_text_box(slide, PInches(0.6), proof_y + PInches(0.45), PInches(7.07), PInches(0.85),
                 show(op.get("proof")), size=14, bold=True, color=PPTX_WHITE)

    # Metadata strip
    md_y = PInches(8.45)
    add_filled_rect(slide, PInches(0.4), md_y, PInches(7.47), PInches(2.6), PPTX_WHITE, line_color=PPTX_PURPLE)
    add_text_box(slide, PInches(0.6), md_y + PInches(0.2), PInches(7.07), PInches(0.3),
                 "METADATA", size=10, bold=True, color=PPTX_PURPLE_DARK)

    rows = [
        ("Product", show(product.get("name"))),
        ("Owner (PM)", show((product.get("owner_pm") or {}).get("name"))),
        ("GTM partner", show((product.get("gtm_partner") or {}).get("name"))),
        ("Version", show(product.get("version"), "v0.1")),
        ("Last reviewed", show(product.get("last_reviewed"))),
        ("Anchor Jira epic", show(product.get("anchor_jira_epic"))),
    ]
    for i, (label, val) in enumerate(rows):
        ry = md_y + PInches(0.55 + i * 0.32)
        add_text_box(slide, PInches(0.7), ry, PInches(2.0), PInches(0.3),
                     label, size=10, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(slide, PInches(2.7), ry, PInches(5.0), PInches(0.3),
                     val, size=10, color=PPTX_TEXT_DARK)

    # Footer
    add_text_box(slide, PInches(0.4), PInches(11.25), PInches(7.5), PInches(0.3),
                 "Oolio  /  GTM & Product Enablement", size=8, color=PPTX_TEXT_GREY)

    prs.save(out_path)


# ============================================================
# 02. Supporting Deck (16:9, 12 slides)
# ============================================================

def build_deck(content, product, out_path):
    """Build the supporting deck per the Oolio style guide.

    Twelve slides. Outcome-style titles. Six allowed slide types.
    No italics for emphasis. White-first colour discipline.
    """
    d = content.get("deck") or {}
    titles = (d.get("slide_titles") or {})
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    SW = prs.slide_width
    SH = prs.slide_height
    TOTAL_SLIDES = 12

    def title_for(key, default_title="[TBC]", default_sub=""):
        block = titles.get(key) or {}
        return show(block.get("title", default_title), default_title), \
               show(block.get("subtitle", default_sub), default_sub)

    def base_slide(*, title=None, subtitle=None, page_n=None):
        s = prs.slides.add_slide(blank)
        add_filled_rect(s, 0, 0, SW, SH, PPTX_WHITE)
        # Top bar
        add_filled_rect(s, 0, 0, SW, PInches(0.5), PPTX_PURPLE_DARK)
        add_text_box(s, PInches(0.5), PInches(0.13), PInches(8.0), PInches(0.3),
                     f"OOLIO  /  GTM SUPPORTING DECK  /  {show(product.get('name'))}",
                     size=10, bold=True, color=PPTX_WHITE)
        if page_n is not None:
            add_text_box(s, SW - PInches(2.5), PInches(0.13), PInches(2.0), PInches(0.3),
                         f"{show(product.get('version'), 'v0.1')}  /  Slide {page_n} of {TOTAL_SLIDES}",
                         size=10, color=PPTX_WHITE, align=PP_ALIGN.RIGHT)
        # Bottom rule
        add_filled_rect(s, 0, SH - PInches(0.05), SW, PInches(0.05), PPTX_PURPLE)
        add_text_box(s, PInches(0.5), SH - PInches(0.4), PInches(8.0), PInches(0.3),
                     "Oolio  /  GTM & Product Enablement", size=9, color=PPTX_TEXT_GREY)
        if title:
            add_text_box(s, PInches(0.5), PInches(0.75), PInches(12.3), PInches(0.7),
                         title, size=28, bold=True, color=PPTX_PURPLE_DARK)
        if subtitle:
            # Style guide: no italics for emphasis. Use weight and colour.
            add_text_box(s, PInches(0.5), PInches(1.5), PInches(12.3), PInches(0.4),
                         subtitle, size=12, color=PPTX_TEXT_GREY)
        return s

    def screenshot_placeholder(slide, x, y, w, h, label):
        """Reserved frame for a real product UI screenshot.

        Style guide section 8: real UI only, no mockups. Marked [SCREENSHOT TBC]
        until the actual capture is wired in by the GTM workstream.
        """
        # Outer dashed-style border via layered rectangles.
        add_filled_rect(slide, x, y, w, h, PRGB(0xF9, 0xF8, 0xFC), line_color=PPTX_PURPLE)
        add_text_box(slide, x + PInches(0.2), y + h / 2 - PInches(0.4),
                     w - PInches(0.4), PInches(0.4),
                     "[SCREENSHOT TBC]", size=14, bold=True,
                     color=PPTX_PURPLE_DARK, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, x + PInches(0.2), y + h / 2,
                     w - PInches(0.4), PInches(0.4),
                     label, size=10, color=PPTX_TEXT_GREY, align=PP_ALIGN.CENTER)

    def product_reveal_slide(reveal, page_n):
        """Slide type. Product reveal. Real UI only, one label, outcome paired."""
        title, subtitle = title_for(reveal.get("id", ""), "[TBC]", "")
        s = base_slide(title=title, subtitle=subtitle, page_n=page_n)
        # Left column: outcome and capability text.
        add_text_box(s, PInches(0.5), PInches(2.1), PInches(5.5), PInches(0.4),
                     "OUTCOME", size=10, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(s, PInches(0.5), PInches(2.45), PInches(5.5), PInches(1.2),
                     show(reveal.get("outcome")), size=18, bold=True, color=PPTX_TEXT_DARK)
        add_text_box(s, PInches(0.5), PInches(3.85), PInches(5.5), PInches(0.4),
                     "CAPABILITY", size=10, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(s, PInches(0.5), PInches(4.2), PInches(5.5), PInches(1.0),
                     show(reveal.get("capability")), size=12, color=PPTX_TEXT_DARK)
        # One-label callout pinned at the bottom.
        add_filled_rect(s, PInches(0.5), PInches(6.05), PInches(5.5), PInches(0.85),
                        PPTX_PURPLE_DARK)
        add_text_box(s, PInches(0.7), PInches(6.1), PInches(5.1), PInches(0.75),
                     show(reveal.get("callout")), size=12, bold=True,
                     color=PPTX_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # Right column: screenshot frame.
        screenshot_placeholder(s, PInches(6.4), PInches(2.1), PInches(6.4), PInches(4.8),
                               show(reveal.get("screenshot_path"), "[SCREENSHOT TBC]"))
        return s

    # Slide 1: Cover (no italics, weight and space only).
    s = prs.slides.add_slide(blank)
    add_filled_rect(s, 0, 0, SW, SH, PPTX_PURPLE_DARK)
    add_text_box(s, PInches(0.6), PInches(0.6), PInches(12), PInches(0.4),
                 "OOLIO  /  GTM SUPPORTING DECK", size=12, bold=True, color=PPTX_WHITE)
    add_text_box(s, PInches(0.6), PInches(2.4), PInches(12), PInches(1.2),
                 show(product.get("name")), size=44, bold=True, color=PPTX_WHITE)
    add_text_box(s, PInches(0.6), PInches(3.7), PInches(12), PInches(1.6),
                 show(product.get("tagline_long")), size=20, color=PPTX_ACCENT)
    add_text_box(s, PInches(0.6), PInches(5.6), PInches(6), PInches(0.4),
                 f"PM  /  {show((product.get('owner_pm') or {}).get('name'))}",
                 size=12, color=PPTX_WHITE)
    add_text_box(s, PInches(0.6), PInches(5.95), PInches(6), PInches(0.4),
                 f"GTM partner  /  {show((product.get('gtm_partner') or {}).get('name'))}",
                 size=12, color=PPTX_WHITE)
    add_text_box(s, PInches(0.6), PInches(6.3), PInches(6), PInches(0.4),
                 f"{show(product.get('version'), 'v0.1')}  /  {show(product.get('last_reviewed'))}",
                 size=12, color=PPTX_ACCENT)

    # Slide 2: Operator problem - Insight type.
    title, subtitle = title_for("problem")
    s = base_slide(title=title, subtitle=subtitle, page_n=2)
    pq = (d.get("problem") or {}).get("quote", "[TBC]")
    add_filled_rect(s, PInches(0.5), PInches(2.1), PInches(12.3), PInches(1.4), PPTX_ACCENT)
    add_text_box(s, PInches(0.8), PInches(2.25), PInches(11.7), PInches(1.1),
                 f'"{pq}"', size=22, bold=True, color=PPTX_PURPLE_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
    bullets = (d.get("problem") or {}).get("context_bullets") or []
    for i, b in enumerate(bullets[:3]):
        y = PInches(3.85 + i * 0.95)
        add_filled_rect(s, PInches(0.5), y, PInches(0.15), PInches(0.8), PPTX_PURPLE)
        add_text_box(s, PInches(0.85), y, PInches(11.9), PInches(0.8),
                     show(b), size=16, color=PPTX_TEXT_DARK,
                     anchor=MSO_ANCHOR.MIDDLE)

    # Slide 3: POV - Insight type.
    title, subtitle = title_for("pov")
    s = base_slide(title=title, subtitle=subtitle, page_n=3)
    big = (d.get("pov") or {}).get("big_idea", "[TBC]")
    add_filled_rect(s, PInches(0.5), PInches(2.1), PInches(12.3), PInches(1.6), PPTX_PURPLE_DARK)
    add_text_box(s, PInches(0.8), PInches(2.2), PInches(11.7), PInches(1.4),
                 big, size=24, bold=True, color=PPTX_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, PInches(0.5), PInches(4.0), PInches(12.3), PInches(3.0),
                 show((d.get("pov") or {}).get("supporting")),
                 size=14, color=PPTX_TEXT_DARK)

    # Slide 4: Who it's for - Insight type with two columns.
    title, subtitle = title_for("icp")
    s = base_slide(title=title, subtitle=subtitle, page_n=4)
    segs = (d.get("icp") or {}).get("segments", [])[:3]
    pers = (d.get("icp") or {}).get("personas", [])[:3]
    add_text_box(s, PInches(0.5), PInches(2.0), PInches(6.0), PInches(0.3),
                 "SEGMENTS", size=10, bold=True, color=PPTX_PURPLE_DARK)
    for i, seg in enumerate(segs):
        y = PInches(2.35 + i * 1.55)
        add_filled_rect(s, PInches(0.5), y, PInches(6.0), PInches(1.4), PPTX_ACCENT)
        add_text_box(s, PInches(0.7), y + PInches(0.1), PInches(5.6), PInches(0.4),
                     show(seg.get("name")), size=14, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(s, PInches(0.7), y + PInches(0.55), PInches(5.6), PInches(0.85),
                     show(seg.get("detail")), size=11, color=PPTX_TEXT_DARK)
    add_text_box(s, PInches(6.85), PInches(2.0), PInches(6.0), PInches(0.3),
                 "PERSONAS", size=10, bold=True, color=PPTX_PURPLE_DARK)
    for i, p in enumerate(pers):
        y = PInches(2.35 + i * 1.55)
        add_filled_rect(s, PInches(6.85), y, PInches(6.0), PInches(1.4), PPTX_WHITE,
                        line_color=PPTX_PURPLE)
        add_text_box(s, PInches(7.05), y + PInches(0.1), PInches(5.6), PInches(0.4),
                     show(p.get("name")), size=14, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(s, PInches(7.05), y + PInches(0.55), PInches(5.6), PInches(0.85),
                     show(p.get("detail")), size=11, color=PPTX_TEXT_DARK)

    # Slides 5, 6, 7: Product reveal type, one per view.
    reveals = d.get("product_reveals") or []
    for idx, reveal in enumerate(reveals[:3]):
        product_reveal_slide(reveal, page_n=5 + idx)

    # Slide 8: Where we win - Decision type. Three columns.
    title, subtitle = title_for("competitive")
    s = base_slide(title=title, subtitle=subtitle, page_n=8)
    comp = d.get("competitive") or {}
    cols = [
        ("WHERE WE COMPETE", show(comp.get("where_we_compete")), PPTX_ACCENT, PPTX_PURPLE_DARK),
        ("WHERE WE WIN", show(comp.get("where_we_win")), PPTX_PURPLE_DARK, PPTX_WHITE),
        ("WHERE WE DON'T", show(comp.get("where_we_dont")), PPTX_WHITE, PPTX_TEXT_DARK),
    ]
    cw = PInches(4.2)
    for i, (label, body, fill, txt) in enumerate(cols):
        x = PInches(0.5 + i * 4.3)
        add_filled_rect(s, x, PInches(2.0), cw, PInches(5.15), fill, line_color=PPTX_PURPLE)
        add_text_box(s, x + PInches(0.2), PInches(2.15), cw - PInches(0.4), PInches(0.4),
                     label, size=12, bold=True, color=txt)
        add_text_box(s, x + PInches(0.2), PInches(2.65), cw - PInches(0.4), PInches(4.4),
                     body, size=11, color=txt)

    # Slide 9: Proof type. Stats and case studies.
    title, subtitle = title_for("proof")
    s = base_slide(title=title, subtitle=subtitle, page_n=9)
    stats = (d.get("proof") or {}).get("stats", [])[:3]
    sw = PInches(4.1)
    for i, st in enumerate(stats):
        x = PInches(0.5 + i * 4.25)
        add_filled_rect(s, x, PInches(2.0), sw, PInches(2.0), PPTX_PURPLE_DARK)
        add_text_box(s, x + PInches(0.2), PInches(2.1), sw - PInches(0.4), PInches(1.2),
                     show(st.get("value")), size=36, bold=True, color=PPTX_WHITE,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, x + PInches(0.2), PInches(3.2), sw - PInches(0.4), PInches(0.7),
                     show(st.get("label")), size=11, color=PPTX_ACCENT)
    cs = (d.get("proof") or {}).get("case_studies", [])[:2]
    for i, c in enumerate(cs):
        y = PInches(4.25 + i * 1.45)
        add_filled_rect(s, PInches(0.5), y, PInches(12.35), PInches(1.35), PPTX_PURPLE_LIGHT,
                        line_color=PPTX_PURPLE)
        add_text_box(s, PInches(0.7), y + PInches(0.1), PInches(11.9), PInches(0.4),
                     show(c.get("customer")), size=12, bold=True, color=PPTX_PURPLE_DARK)
        add_text_box(s, PInches(0.7), y + PInches(0.55), PInches(11.9), PInches(0.75),
                     show(c.get("story")), size=11, color=PPTX_TEXT_DARK)

    # Slide 10: How to buy - Process type. Five horizontal steps.
    title, subtitle = title_for("buy_steps")
    s = base_slide(title=title, subtitle=subtitle, page_n=10)
    steps = d.get("buy_steps") or []
    sw = PInches(2.5)
    for i, st in enumerate(steps[:5]):
        x = PInches(0.5 + i * 2.55)
        add_filled_rect(s, x, PInches(2.1), sw, PInches(0.7), PPTX_PURPLE_DARK)
        add_text_box(s, x + PInches(0.1), PInches(2.2), sw - PInches(0.2), PInches(0.5),
                     f"{i+1}. {show(st.get('name'))}", size=12, bold=True, color=PPTX_WHITE,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_filled_rect(s, x, PInches(2.8), sw, PInches(4.2), PPTX_WHITE, line_color=PPTX_PURPLE)
        add_text_box(s, x + PInches(0.15), PInches(2.95), sw - PInches(0.3), PInches(3.9),
                     show(st.get("description")), size=10, color=PPTX_TEXT_DARK)

    # Slide 11: Summary - Decision type. Takeaway plus next steps with owner and date.
    title, subtitle = title_for("summary")
    s = base_slide(title=title, subtitle=subtitle, page_n=11)
    summary = d.get("summary") or {}
    add_filled_rect(s, PInches(0.5), PInches(2.1), PInches(12.35), PInches(1.5), PPTX_ACCENT)
    add_text_box(s, PInches(0.7), PInches(2.2), PInches(12.0), PInches(1.3),
                 show(summary.get("takeaway")), size=14, bold=True, color=PPTX_PURPLE_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
    headers = ["Action", "Owner", "Due"]
    col_widths = [PInches(7.5), PInches(2.5), PInches(2.35)]
    col_x = [PInches(0.5), PInches(8.0), PInches(10.5)]
    for i, h in enumerate(headers):
        add_filled_rect(s, col_x[i], PInches(3.85), col_widths[i], PInches(0.45), PPTX_PURPLE_DARK)
        add_text_box(s, col_x[i] + PInches(0.15), PInches(3.92), col_widths[i] - PInches(0.3), PInches(0.32),
                     h, size=11, bold=True, color=PPTX_WHITE)
    rows = (summary.get("next_steps") or [])[:3]
    for r_i, row in enumerate(rows):
        y = PInches(4.35 + r_i * 0.85)
        zebra = (r_i % 2 == 0)
        fill = PPTX_PURPLE_LIGHT if zebra else PPTX_WHITE
        vals = [show(row.get("action")), show(row.get("owner")), show(row.get("due"))]
        for c_i, val in enumerate(vals):
            add_filled_rect(s, col_x[c_i], y, col_widths[c_i], PInches(0.8), fill,
                            line_color=PRGB(0xD5, 0xD0, 0xE1))
            add_text_box(s, col_x[c_i] + PInches(0.15), y + PInches(0.1),
                         col_widths[c_i] - PInches(0.3), PInches(0.65),
                         val, size=10, color=PPTX_TEXT_DARK)

    # Slide 12: Back cover.
    s = prs.slides.add_slide(blank)
    add_filled_rect(s, 0, 0, SW, SH, PPTX_PURPLE_DARK)
    add_text_box(s, PInches(0.6), PInches(0.6), PInches(12), PInches(0.4),
                 "OOLIO  /  THANK YOU", size=12, bold=True, color=PPTX_WHITE)
    promise = show((d.get("back_cover") or {}).get("promise"))
    add_text_box(s, PInches(0.6), PInches(2.4), PInches(12), PInches(1.6),
                 promise, size=36, bold=True, color=PPTX_WHITE)
    add_text_box(s, PInches(0.6), PInches(4.2), PInches(12), PInches(0.4),
                 "CONTACT", size=12, bold=True, color=PPTX_ACCENT)
    contacts = (d.get("back_cover") or {}).get("contacts") or []
    for i, c in enumerate(contacts):
        y = PInches(4.7 + i * 0.6)
        add_text_box(s, PInches(0.6), y, PInches(3.0), PInches(0.4),
                     show(c.get("role")), size=11, bold=True, color=PPTX_ACCENT)
        add_text_box(s, PInches(3.7), y, PInches(3.5), PInches(0.4),
                     show(c.get("name")), size=11, color=PPTX_WHITE)
        add_text_box(s, PInches(7.4), y, PInches(5.5), PInches(0.4),
                     show(c.get("email")), size=11, color=PPTX_ACCENT)

    prs.save(out_path)


# ============================================================
# Main
# ============================================================

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 build_pack.py <path to pack_content.json>", file=sys.stderr)
        sys.exit(2)
    json_path = sys.argv[1]
    with open(json_path, "r", encoding="utf-8") as f:
        content = json.load(f)
    out_dir = os.path.dirname(os.path.abspath(json_path))
    product = content.get("product") or {}
    version = show(product.get("version"), "v0.1")

    targets = [
        (f"01_One_Pager_{version}.pptx", lambda p: build_one_pager(content, product, p)),
        (f"02_Supporting_Deck_{version}.pptx", lambda p: build_deck(content, product, p)),
        (f"03_Sales_Playbook_{version}.docx", lambda p: build_sales(content, product, p)),
        (f"04_Account_Management_Playbook_{version}.docx", lambda p: build_am(content, product, p)),
        (f"05_Onboarding_Playbook_{version}.docx", lambda p: build_onboarding(content, product, p)),
        (f"06_Marketing_Pack_{version}.docx", lambda p: build_marketing(content, product, p)),
    ]

    print(f"Building {len(targets)} files into {out_dir} ...")
    for name, fn in targets:
        path = os.path.join(out_dir, name)
        fn(path)
        print(f"  Wrote {name}")
    print("Done.")


if __name__ == "__main__":
    main()
